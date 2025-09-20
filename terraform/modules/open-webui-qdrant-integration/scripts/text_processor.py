#!/usr/bin/env python3
"""
Text Processor for Open WebUI Qdrant Integration

This module provides text processing capabilities including text extraction,
chunking, keyword extraction, and preprocessing for document indexing.
"""

import os
import re
import logging
from typing import Dict, List, Optional, Any, Union, Tuple
import mimetypes
from io import BytesIO
import asyncio

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TextProcessor:
    def __init__(self):
        self.chunk_size = int(os.getenv('TEXT_CHUNK_SIZE', '1000'))
        self.chunk_overlap = int(os.getenv('TEXT_CHUNK_OVERLAP', '200'))
        self.max_keywords = int(os.getenv('MAX_KEYWORDS', '20'))
        
        # Initialize text extraction libraries
        self._initialize_extractors()
    
    def _initialize_extractors(self):
        """Initialize text extraction libraries"""
        try:
            # Try to import optional dependencies
            global PyPDF2, docx, openpyxl, python_pptx, BeautifulSoup
            
            try:
                import PyPDF2
                self.pdf_available = True
            except ImportError:
                logger.warning("PyPDF2 not available - PDF extraction disabled")
                self.pdf_available = False
            
            try:
                import docx
                self.docx_available = True
            except ImportError:
                logger.warning("python-docx not available - DOCX extraction disabled")
                self.docx_available = False
            
            try:
                import openpyxl
                self.excel_available = True
            except ImportError:
                logger.warning("openpyxl not available - Excel extraction disabled")
                self.excel_available = False
            
            try:
                from pptx import Presentation
                global Presentation
                self.pptx_available = True
            except ImportError:
                logger.warning("python-pptx not available - PowerPoint extraction disabled")
                self.pptx_available = False
            
            try:
                from bs4 import BeautifulSoup
                self.html_available = True
            except ImportError:
                logger.warning("BeautifulSoup not available - HTML extraction disabled")
                self.html_available = False
            
            logger.info("Text extraction libraries initialized")
            
        except Exception as e:
            logger.error(f"Failed to initialize text extractors: {e}")
    
    async def extract_text(self, content: bytes, content_type: str, filename: str) -> str:
        """Extract text from document content based on content type"""
        try:
            # Determine content type if not provided
            if not content_type:
                content_type, _ = mimetypes.guess_type(filename)
                content_type = content_type or 'application/octet-stream'
            
            logger.info(f"Extracting text from {content_type} document: {filename}")
            
            # Route to appropriate extractor
            if content_type == 'application/pdf':
                return await self._extract_pdf_text(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
                return await self._extract_docx_text(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
                return await self._extract_excel_text(content)
            elif content_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
                return await self._extract_pptx_text(content)
            elif content_type in ['text/html', 'application/xhtml+xml']:
                return await self._extract_html_text(content)
            elif content_type.startswith('text/'):
                return await self._extract_plain_text(content)
            else:
                # Try to extract as plain text
                logger.warning(f"Unknown content type {content_type}, attempting plain text extraction")
                return await self._extract_plain_text(content)
                
        except Exception as e:
            logger.error(f"Failed to extract text from {filename}: {e}")
            # Fallback to plain text extraction
            try:
                return await self._extract_plain_text(content)
            except:
                return ""
    
    async def _extract_pdf_text(self, content: bytes) -> str:
        """Extract text from PDF content"""
        if not self.pdf_available:
            raise ValueError("PDF extraction not available - PyPDF2 not installed")
        
        try:
            text_parts = []
            pdf_file = BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(page_text)
                except Exception as e:
                    logger.warning(f"Failed to extract text from PDF page {page_num}: {e}")
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"PDF text extraction failed: {e}")
            raise
    
    async def _extract_docx_text(self, content: bytes) -> str:
        """Extract text from DOCX content"""
        if not self.docx_available:
            raise ValueError("DOCX extraction not available - python-docx not installed")
        
        try:
            doc_file = BytesIO(content)
            doc = docx.Document(doc_file)
            
            text_parts = []
            
            # Extract paragraph text
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract table text
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"DOCX text extraction failed: {e}")
            raise
    
    async def _extract_excel_text(self, content: bytes) -> str:
        """Extract text from Excel content"""
        if not self.excel_available:
            raise ValueError("Excel extraction not available - openpyxl not installed")
        
        try:
            excel_file = BytesIO(content)
            workbook = openpyxl.load_workbook(excel_file, data_only=True)
            
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value))
                    if row_text:
                        text_parts.append(' | '.join(row_text))
                
                text_parts.append("")  # Add spacing between sheets
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Excel text extraction failed: {e}")
            raise
    
    async def _extract_pptx_text(self, content: bytes) -> str:
        """Extract text from PowerPoint content"""
        if not self.pptx_available:
            raise ValueError("PowerPoint extraction not available - python-pptx not installed")
        
        try:
            pptx_file = BytesIO(content)
            presentation = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = [f"Slide {slide_num + 1}:"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the slide number
                    text_parts.append('\n'.join(slide_text))
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"PowerPoint text extraction failed: {e}")
            raise
    
    async def _extract_html_text(self, content: bytes) -> str:
        """Extract text from HTML content"""
        if not self.html_available:
            raise ValueError("HTML extraction not available - BeautifulSoup not installed")
        
        try:
            html_content = content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and clean it up
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text
            
        except Exception as e:
            logger.error(f"HTML text extraction failed: {e}")
            raise
    
    async def _extract_plain_text(self, content: bytes) -> str:
        """Extract plain text content"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    text = content.decode(encoding)
                    return text
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, use utf-8 with error handling
            return content.decode('utf-8', errors='ignore')
            
        except Exception as e:
            logger.error(f"Plain text extraction failed: {e}")
            raise
    
    async def chunk_text(self, text: str, chunk_size: Optional[int] = None, overlap: Optional[int] = None) -> List[str]:
        """Split text into overlapping chunks"""
        try:
            chunk_size = chunk_size or self.chunk_size
            overlap = overlap or self.chunk_overlap
            
            if len(text) <= chunk_size:
                return [text]
            
            chunks = []
            start = 0
            
            while start < len(text):
                # Calculate end position
                end = start + chunk_size
                
                # If this is not the last chunk, try to break at a sentence or word boundary
                if end < len(text):
                    # Look for sentence boundaries within the last 200 characters
                    search_start = max(end - 200, start)
                    sentence_end = self._find_sentence_boundary(text, search_start, end)
                    
                    if sentence_end > start:
                        end = sentence_end
                    else:
                        # Look for word boundaries
                        word_end = self._find_word_boundary(text, end)
                        if word_end > start:
                            end = word_end
                
                # Extract chunk
                chunk = text[start:end].strip()
                if chunk:
                    chunks.append(chunk)
                
                # Move start position with overlap
                start = max(start + 1, end - overlap)
                
                # Prevent infinite loop
                if start >= len(text):
                    break
            
            logger.debug(f"Split text into {len(chunks)} chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Text chunking failed: {e}")
            return [text]  # Return original text as single chunk
    
    def _find_sentence_boundary(self, text: str, start: int, end: int) -> int:
        """Find the best sentence boundary within the range"""
        # Look for sentence endings
        sentence_endings = ['. ', '! ', '? ', '.\n', '!\n', '?\n']
        
        best_pos = start
        for i in range(end - 1, start - 1, -1):
            for ending in sentence_endings:
                if text[i:i + len(ending)] == ending:
                    return i + 1
        
        return best_pos
    
    def _find_word_boundary(self, text: str, pos: int) -> int:
        """Find the nearest word boundary before the position"""
        # Look backwards for whitespace
        for i in range(pos - 1, max(0, pos - 100), -1):
            if text[i].isspace():
                return i
        
        return pos
    
    async def extract_keywords(self, text: str, max_keywords: Optional[int] = None) -> List[str]:
        """Extract keywords from text using simple frequency analysis"""
        try:
            max_keywords = max_keywords or self.max_keywords
            
            # Clean and normalize text
            cleaned_text = self._clean_text_for_keywords(text)
            
            # Split into words
            words = cleaned_text.split()
            
            # Filter words
            filtered_words = []
            for word in words:
                if (len(word) >= 3 and 
                    word.isalpha() and 
                    word.lower() not in self._get_stopwords()):
                    filtered_words.append(word.lower())
            
            # Count word frequencies
            word_freq = {}
            for word in filtered_words:
                word_freq[word] = word_freq.get(word, 0) + 1
            
            # Sort by frequency and return top keywords
            sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
            keywords = [word for word, freq in sorted_words[:max_keywords]]
            
            logger.debug(f"Extracted {len(keywords)} keywords from text")
            return keywords
            
        except Exception as e:
            logger.error(f"Keyword extraction failed: {e}")
            return []
    
    def _clean_text_for_keywords(self, text: str) -> str:
        """Clean text for keyword extraction"""
        # Remove special characters and normalize whitespace
        text = re.sub(r'[^\w\s]', ' ', text)
        text = re.sub(r'\s+', ' ', text)
        return text.strip()
    
    def _get_stopwords(self) -> set:
        """Get common English stopwords"""
        return {
            'a', 'an', 'and', 'are', 'as', 'at', 'be', 'by', 'for', 'from',
            'has', 'he', 'in', 'is', 'it', 'its', 'of', 'on', 'that', 'the',
            'to', 'was', 'will', 'with', 'would', 'you', 'your', 'have', 'had',
            'this', 'these', 'they', 'were', 'been', 'their', 'said', 'each',
            'which', 'she', 'do', 'how', 'if', 'we', 'when', 'where', 'who',
            'why', 'what', 'can', 'could', 'should', 'may', 'might', 'must',
            'shall', 'about', 'after', 'all', 'also', 'am', 'another', 'any',
            'because', 'before', 'being', 'between', 'both', 'but', 'came',
            'come', 'did', 'does', 'during', 'each', 'few', 'get', 'got',
            'here', 'him', 'his', 'into', 'just', 'like', 'make', 'many',
            'me', 'more', 'most', 'much', 'my', 'new', 'no', 'not', 'now',
            'only', 'or', 'other', 'our', 'out', 'over', 'own', 'see', 'so',
            'some', 'such', 'than', 'them', 'there', 'through', 'time', 'two',
            'up', 'use', 'used', 'using', 'very', 'way', 'well', 'work'
        }
    
    async def preprocess_text(self, text: str) -> str:
        """Preprocess text for better embedding generation"""
        try:
            # Remove excessive whitespace
            text = re.sub(r'\s+', ' ', text)
            
            # Remove control characters
            text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)
            
            # Normalize quotes and dashes
            text = text.replace('"', '"').replace('"', '"')
            text = text.replace(''', "'").replace(''', "'")
            text = text.replace('–', '-').replace('—', '-')
            
            # Remove excessive punctuation
            text = re.sub(r'[.]{3,}', '...', text)
            text = re.sub(r'[!]{2,}', '!', text)
            text = re.sub(r'[?]{2,}', '?', text)
            
            # Clean up spacing around punctuation
            text = re.sub(r'\s+([,.!?;:])', r'\1', text)
            text = re.sub(r'([,.!?;:])\s+', r'\1 ', text)
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"Text preprocessing failed: {e}")
            return text
    
    async def get_text_stats(self, text: str) -> Dict[str, Any]:
        """Get statistics about the text"""
        try:
            words = text.split()
            sentences = re.split(r'[.!?]+', text)
            paragraphs = text.split('\n\n')
            
            return {
                'character_count': len(text),
                'word_count': len(words),
                'sentence_count': len([s for s in sentences if s.strip()]),
                'paragraph_count': len([p for p in paragraphs if p.strip()]),
                'average_word_length': sum(len(word) for word in words) / len(words) if words else 0,
                'average_sentence_length': len(words) / len(sentences) if sentences else 0
            }
            
        except Exception as e:
            logger.error(f"Failed to get text statistics: {e}")
            return {}
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported document formats"""
        formats = ['text/plain', 'text/html']
        
        if self.pdf_available:
            formats.append('application/pdf')
        
        if self.docx_available:
            formats.extend([
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'application/msword'
            ])
        
        if self.excel_available:
            formats.extend([
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'application/vnd.ms-excel'
            ])
        
        if self.pptx_available:
            formats.extend([
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'application/vnd.ms-powerpoint'
            ])
        
        return formats
    
    async def validate_document(self, content: bytes, content_type: str, filename: str) -> Dict[str, Any]:
        """Validate document and return information"""
        try:
            # Check file size
            size_mb = len(content) / (1024 * 1024)
            max_size_mb = 50  # 50MB limit
            
            if size_mb > max_size_mb:
                return {
                    'valid': False,
                    'error': f'Document too large: {size_mb:.1f}MB (max: {max_size_mb}MB)'
                }
            
            # Check content type
            supported_formats = self.get_supported_formats()
            if content_type not in supported_formats:
                return {
                    'valid': False,
                    'error': f'Unsupported format: {content_type}'
                }
            
            # Try to extract a small sample of text
            try:
                sample_text = await self.extract_text(content[:10000], content_type, filename)
                if not sample_text.strip():
                    return {
                        'valid': False,
                        'error': 'No text could be extracted from document'
                    }
            except Exception as e:
                return {
                    'valid': False,
                    'error': f'Text extraction failed: {str(e)}'
                }
            
            return {
                'valid': True,
                'size_mb': size_mb,
                'content_type': content_type,
                'extractable': True
            }
            
        except Exception as e:
            logger.error(f"Document validation failed: {e}")
            return {
                'valid': False,
                'error': f'Validation failed: {str(e)}'
            }